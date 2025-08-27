# loaders.py
from __future__ import annotations
import os, re
from pathlib import Path
from typing import Iterable, List

import chardet
import fitz  # PyMuPDF
from pypdf import PdfReader
import docx2txt
from pptx import Presentation
import pandas as pd

TEXT_EXTS = {".txt", ".md", ".csv"}
DOC_EXTS  = {".docx"}
PPT_EXTS  = {".pptx"}
XLS_EXTS  = {".xlsx"}
PDF_EXTS  = {".pdf"}
SUPPORTED = TEXT_EXTS | DOC_EXTS | PPT_EXTS | XLS_EXTS | PDF_EXTS

_ws = re.compile(r"[ \t\u00A0]+")

# Safety cap for huge PDFs
MAX_PDF_PAGES = int(os.getenv("RAG_MAX_PDF_PAGES", "200"))

def _clean(s: str) -> str:
    s = s.replace("\x00", " ")
    s = re.sub(r"\r\n?|\n", "\n", s)
    s = _ws.sub(" ", s)
    return s.strip()

def from_text_like(path: Path) -> List[str]:
    raw = path.read_bytes()
    enc = chardet.detect(raw).get("encoding") or "utf-8"
    try:
        txt = raw.decode(enc, errors="ignore")
    except Exception:
        txt = raw.decode("utf-8", errors="ignore")
    return [_clean(txt)]

def from_pdf_pymupdf(path: Path) -> List[str]:
    out: List[str] = []
    with fitz.open(str(path)) as doc:
        n = min(len(doc), MAX_PDF_PAGES)
        for i in range(n):
            t = doc[i].get_text("text") or ""
            if t.strip():
                out.append(_clean(t))
    return out

def from_pdf_pypdf(path: Path) -> List[str]:
    out: List[str] = []
    r = PdfReader(str(path))
    n = min(len(r.pages), MAX_PDF_PAGES)
    for i in range(n):
        t = (r.pages[i].extract_text() or "").strip()
        if t:
            out.append(_clean(t))
    return out

def from_pdf(path: Path) -> List[str]:
    # 1) try PyMuPDF (fast)
    try:
        return from_pdf_pymupdf(path)
    except Exception:
        pass
    # 2) fallback to pypdf (slower, but sometimes useful)
    try:
        return from_pdf_pypdf(path)
    except Exception:
        return []

def from_docx(path: Path) -> List[str]:
    text = docx2txt.process(str(path)) or ""
    return [_clean(text)] if text else []

def from_pptx(path: Path) -> List[str]:
    prs = Presentation(str(path))
    buf: List[str] = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shp in slide.shapes:
            if hasattr(shp, "has_text_frame") and shp.has_text_frame:
                parts.append(shp.text)
        if parts:
            buf.append(_clean(f"Slide {i+1}: " + " \n ".join(parts)))
    return buf

def from_xlsx(path: Path, max_cells: int = 5000) -> List[str]:
    dfs = pd.read_excel(str(path), sheet_name=None, dtype=str)
    out: List[str] = []
    cells = 0
    for name, df in dfs.items():
        df = df.fillna("")
        for _, row in df.iterrows():
            if cells > max_cells:
                break
            line = " | ".join(map(str, row.tolist()))
            if line.strip():
                out.append(_clean(f"{name}: {line}"))
                cells += len(row)
        if cells > max_cells:
            break
    return out

def load_file(path: Path) -> List[str]:
    ext = path.suffix.lower()
    if ext in PDF_EXTS:  return from_pdf(path)
    if ext in DOC_EXTS:  return from_docx(path)
    if ext in PPT_EXTS:  return from_pptx(path)
    if ext in XLS_EXTS:  return from_xlsx(path)
    if ext in TEXT_EXTS: return from_text_like(path)
    return []

def iter_files(folder: Path) -> Iterable[Path]:
    # Allow temporarily disabling PDFs for debugging:
    disable_pdf = os.getenv("RAG_DISABLE_PDF", "0") == "1"
    for p in folder.rglob("*"):
        if p.is_file():
            ext = p.suffix.lower()
            if disable_pdf and ext in PDF_EXTS:
                continue
            if ext in SUPPORTED:
                yield p
